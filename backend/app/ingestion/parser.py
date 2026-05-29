# backend/app/ingestion/parser.py
import hashlib
import logging
import os
from datetime import datetime, timezone
from email import policy
from email.parser import BytesParser
from typing import Optional

from app.models.source import DocumentSection, ParsedDocument

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".docx", ".pptx", ".xlsx", ".pdf", ".msg", ".eml", ".txt", ".md"}


def parse_document(file_path: str) -> ParsedDocument:
    ext = os.path.splitext(file_path)[1].lower()
    last_modified = datetime.fromtimestamp(os.path.getmtime(file_path), tz=timezone.utc)

    parsers = {
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".pdf": _parse_pdf,
        ".msg": _parse_msg,
        ".eml": _parse_eml,
        ".txt": _parse_text,
        ".md": _parse_text,
    }

    parser_fn = parsers.get(ext)
    if parser_fn is None:
        raise ValueError(f"Unsupported file type: {ext}")

    sections = parser_fn(file_path)

    return ParsedDocument(
        file_path=file_path,
        file_type=ext.lstrip("."),
        last_modified=last_modified,
        sections=sections,
        metadata={"size_bytes": os.path.getsize(file_path)},
    )


def _parse_docx(path: str) -> list[DocumentSection]:
    from docx import Document
    doc = Document(path)
    sections = []
    current_title = None
    current_text = []

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith("Heading"):
            if current_text:
                sections.append(DocumentSection(
                    title=current_title, text="\n".join(current_text)
                ))
                current_text = []
            current_title = para.text
        else:
            if para.text.strip():
                current_text.append(para.text)

    if current_text:
        sections.append(DocumentSection(title=current_title, text="\n".join(current_text)))
    if not sections:
        full_text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        sections.append(DocumentSection(title=None, text=full_text))
    return sections


def _parse_pptx(path: str) -> list[DocumentSection]:
    from pptx import Presentation
    prs = Presentation(path)
    sections = []
    for i, slide in enumerate(prs.slides, 1):
        title = None
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape == slide.shapes.title:
                    title = shape.text_frame.text
                else:
                    text = shape.text_frame.text.strip()
                    if text:
                        texts.append(text)
        if title or texts:
            sections.append(DocumentSection(
                title=title, text="\n".join(texts), page_number=i
            ))
    return sections


def _parse_xlsx(path: str) -> list[DocumentSection]:
    from openpyxl import load_workbook
    wb = load_workbook(path, read_only=True, data_only=True)
    sections = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        rows = []
        headers = []
        for i, row in enumerate(ws.iter_rows(values_only=True)):
            values = [str(c) if c is not None else "" for c in row]
            if i == 0:
                headers = values
            else:
                rows.append(" | ".join(values))
        text = ""
        if headers:
            text = "Columns: " + " | ".join(headers) + "\n"
        text += "\n".join(rows)
        if text.strip():
            sections.append(DocumentSection(title=sheet, text=text))
    wb.close()
    return sections


def _parse_pdf(path: str) -> list[DocumentSection]:
    sections = []
    try:
        import pdfplumber
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    sections.append(DocumentSection(
                        title=None, text=text, page_number=i
                    ))
    except Exception:
        logger.info("pdfplumber failed for %s, trying pymupdf", path)
        import fitz
        doc = fitz.open(path)
        for i, page in enumerate(doc, 1):
            text = page.get_text()
            if text.strip():
                sections.append(DocumentSection(
                    title=None, text=text, page_number=i
                ))
        doc.close()
    return sections


def _parse_msg(path: str) -> list[DocumentSection]:
    import extract_msg
    msg = extract_msg.Message(path)
    subject = msg.subject or "No Subject"
    sender = msg.sender or "Unknown"
    body = msg.body or ""
    text = f"From: {sender}\nSubject: {subject}\n\n{body}"
    msg.close()
    return [DocumentSection(title=subject, text=text)]


def _parse_eml(path: str) -> list[DocumentSection]:
    with open(path, "rb") as f:
        msg = BytesParser(policy=policy.default).parse(f)
    subject = msg.get("subject", "No Subject")
    sender = msg.get("from", "Unknown")
    body = msg.get_body(preferencelist=("plain", "html"))
    text_content = body.get_content() if body else ""
    text = f"From: {sender}\nSubject: {subject}\n\n{text_content}"
    return [DocumentSection(title=subject, text=text)]


def _parse_text(path: str) -> list[DocumentSection]:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    return [DocumentSection(title=None, text=content)]
