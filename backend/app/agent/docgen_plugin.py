# app/agent/docgen_plugin.py
import json
import os
from datetime import datetime, timezone
from semantic_kernel.functions import kernel_function
from app.config import get_settings


class DocumentGenerationPlugin:
    @kernel_function(
        name="generate_presentation",
        description="Generate a PowerPoint presentation. Provide a title and slides as JSON array: [{\"title\": \"...\", \"bullets\": [\"...\"]}]"
    )
    async def generate_pptx(self, title: str, slides_json: str, client_name: str = "default") -> str:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        slides_data = json.loads(slides_json)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for slide_data in slides_data:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_data.get("title", "")
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for bullet in slide_data.get("bullets", []):
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

        settings = get_settings()
        out_dir = os.path.join(settings.ONEDRIVE_SYNC_PATH, client_name, "Generated")
        os.makedirs(out_dir, exist_ok=True)
        ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        filename = f"{ts}_{title.replace(' ', '_')}.pptx"
        out_path = os.path.join(out_dir, filename)
        prs.save(out_path)

        return json.dumps({"status": "created", "file_path": os.path.relpath(out_path, settings.ONEDRIVE_SYNC_PATH)})

    @kernel_function(
        name="generate_document",
        description="Generate a Word document. Provide a title and sections as JSON array: [{\"heading\": \"...\", \"content\": \"...\"}]"
    )
    async def generate_docx(self, title: str, sections_json: str, client_name: str = "default") -> str:
        from docx import Document

        sections_data = json.loads(sections_json)
        doc = Document()
        doc.add_heading(title, level=0)

        for section in sections_data:
            doc.add_heading(section.get("heading", ""), level=1)
            doc.add_paragraph(section.get("content", ""))

        settings = get_settings()
        out_dir = os.path.join(settings.ONEDRIVE_SYNC_PATH, client_name, "Generated")
        os.makedirs(out_dir, exist_ok=True)
        ts = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
        filename = f"{ts}_{title.replace(' ', '_')}.docx"
        out_path = os.path.join(out_dir, filename)
        doc.save(out_path)

        return json.dumps({"status": "created", "file_path": os.path.relpath(out_path, settings.ONEDRIVE_SYNC_PATH)})
